import streamlit as st
from pptx import Presentation
from pptx.util import Cm  # Using centimeters for measurements

def modify_boxes():
    # Load the presentation
    prs = Presentation('plakat_pp.pptx')
    slide = prs.slides[0]
    
    # Original and new dimensions
    original_width = Cm(6)    # 6 cm
    original_height = Cm(5.5) # 5.5 cm
    new_width = Cm(12)       # 12 cm
    new_height = Cm(11)      # 11 cm
    gap = Cm(0.2)           # 2 mm gap
    
    # Find all boxes that match the original dimensions (with some tolerance)
    tolerance = Cm(0.1)  # 1mm tolerance for dimension matching
    boxes_to_resize = []
    
    for shape in slide.shapes:
        # Debug: Print dimensions of each shape
        st.write(f"Shape dimensions: Width={shape.width/Cm(1):.2f}cm, Height={shape.height/Cm(1):.2f}cm")
        
        if abs(shape.width - original_width) <= tolerance and abs(shape.height - original_height) <= tolerance:
            boxes_to_resize.append(shape)
            st.write(f"Found box to resize: {shape.text[:30] if hasattr(shape, 'text_frame') else 'No text'}")
    
    if not boxes_to_resize:
        st.warning("No boxes with 6x5.5cm dimensions found!")
        return
    
    # Calculate scale factors and gap
    width_scale = new_width / original_width
    height_scale = new_height / original_height
    gap = Cm(0.2)  # 2mm gap
    
    # First, store all original positions and sort boxes by position
    box_positions = [(shape, shape.left, shape.top) for shape in boxes_to_resize]
    # Sort by top position first, then by left position
    box_positions.sort(key=lambda x: (x[2], x[1]))
    
    # Track the last box's position to maintain gaps
    last_left = None
    current_row_top = None
    previous_row_top = None
    
    # Resize and reposition boxes
    for shape, original_left, original_top in box_positions:
        # Resize the box
        shape.width = new_width
        shape.height = new_height
        
        if current_row_top is None or abs(original_top - current_row_top) > original_height / 2:
            # This is either the first box or start of a new row
            if current_row_top is not None:
                # If this is a new row (not the first box), position it below the previous row with gap
                previous_row_top = current_row_top
                shape.top = round(previous_row_top + new_height + gap)
            else:
                # First box in the presentation
                shape.top = round(original_top - (new_height - original_height) / 2)
            
            shape.left = round(original_left - (new_width - original_width) / 2)
            last_left = shape.left
            current_row_top = original_top
        else:
            # Same row, position to the right of previous box
            shape.left = round(last_left + new_width + gap)
            shape.top = round(previous_row_top + new_height + gap if previous_row_top is not None else original_top)
            last_left = shape.left
    
    # Save the modified presentation
    output_path = "modified_plakat.pptx"
    prs.save(output_path)
    
    # Create download button
    with open(output_path, "rb") as file:
        st.download_button(
            label="Download Modified PowerPoint",
            data=file,
            file_name="modified_plakat.pptx",
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
        )
    
    st.success(f"Found and resized {len(boxes_to_resize)} boxes to 12x11 cm!")

# Streamlit interface
st.title('PowerPoint Box Modifier')
st.write("This will modify only the 6x5.5cm boxes to 12x11 cm while maintaining their relative positions.")

if st.button('Modify Boxes'):
    modify_boxes()
